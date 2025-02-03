import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

# Define the file path
file_path = "Student_Result_Management_System_Presentation.pptx"

# Check if the file exists
if not os.path.exists(file_path):
    print(f"Error: File not found at {file_path}")
else:
    # Load the existing presentation
    presentation = Presentation(file_path)

    # Function to apply a theme to the presentation
    def apply_theme(presentation):
        for slide in presentation.slides:
            for shape in slide.shapes:
                # Check if the shape has text (to avoid errors on non-text shapes)
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        # Style for Title Placeholders
                        if shape.placeholder_format.idx == 0:  # Title placeholder
                            paragraph.font.size = Pt(36)
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = RGBColor(0, 51, 102)  # Dark Blue
                        # Style for Content Placeholders
                        else:  # Content placeholder
                            paragraph.font.size = Pt(20)
                            paragraph.font.color.rgb = RGBColor(0, 102, 204)  # Light Blue

    # Apply the theme to the presentation
    apply_theme(presentation)

    # Save the styled presentation
    new_file_path = "Student_Result_Management_System_Styled_Presentation.pptx"
    presentation.save(new_file_path)
    print(f"Styled presentation saved as {new_file_path}")
