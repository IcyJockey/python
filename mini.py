from pptx import Presentation

# Create a PowerPoint presentation
presentation = Presentation()

# Slide 1: Title Slide
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Student Result Management System"
subtitle.text = (
    "A Project Report\n\nPresented By: [Your Name(s)]\nRoll Number: [Your Roll Number(s)]\n"
    "Under the Guidance of: [Your Mentor's Name]\n[Institution Name]"
)

# Slide 2: Introduction
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Introduction"
content.text = (
    "• Overview of the project.\n"
    "• Purpose: Simplify the process of managing student results.\n"
    "• Key users: Teachers and students."
)

# Slide 3: Objectives
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Objectives"
content.text = (
    "• Provide an easy-to-use interface for managing student data.\n"
    "• Enable secure access for teachers and students.\n"
    "• Automate grade calculation based on marks."
)

# Slide 4: Features
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Features"
content.text = (
    "Teacher Features:\n"
    "  • Add, update, view, delete results.\n"
    "  • Search student records.\n\n"
    "Student Features:\n"
    "  • Secure login.\n"
    "  • View results.\n\n"
    "Automatic grade calculation."
)

# Slide 5: Technologies Used
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Technologies Used"
content.text = (
    "• Programming Language: Python\n"
    "• Database: SQLite\n"
    "• GUI Library: Tkinter\n"
    "• IDE: IDLE/VS Code"
)

# Slide 6: Implementation
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Implementation"
content.text = (
    "• Database Design: One table for storing student data.\n"
    "• User Interface:\n"
    "  - Login screens for teachers and students.\n"
    "  - Dashboards for managing and viewing data.\n"
    "• Key Features:\n"
    "  - Data validation.\n"
    "  - Error handling (e.g., duplicate roll numbers)."
)

# Slide 7: Results
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Results"
content.text = (
    "• Successfully added and updated student records.\n"
    "• Secure access for teachers and students.\n"
    "• User-friendly interface with clear navigation."
)

# Slide 8: Conclusion
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]
title.text = "Conclusion"
content.text = (
    "• Efficient tool for managing academic records.\n"
    "• Streamlines operations for teachers and students.\n"
    "• Future scope: Adding Excel export and role-based permissions."
)

# Save the presentation
presentation.save("Student_Result_Management_System_Presentation.pptx")
